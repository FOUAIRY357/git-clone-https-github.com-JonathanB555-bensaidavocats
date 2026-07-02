#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Generateur de support de cadrage strategique - charte BENSAID AVOCATS."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import sys

GOLD=RGBColor(0xB0,0x8A,0x46); INK=RGBColor(0x11,0x11,0x11); WHITE=RGBColor(0xFF,0xFF,0xFF)
GREYK=RGBColor(0x6E,0x6E,0x6E); GREY9=RGBColor(0x9A,0x9A,0x9A); GREYD=RGBColor(0xD6,0xD6,0xD6)
GREEN=RGBColor(0x2E,0x7D,0x32); RED=RGBColor(0xC0,0x39,0x2B); BLACK=RGBColor(0x00,0x00,0x00)
PANEL=RGBColor(0x16,0x16,0x16); RULEL=RGBColor(0xDA,0xDA,0xDA); RULED=RGBColor(0x33,0x33,0x33)
F="Arial"

prs=Presentation(); prs.slide_width=Emu(12192000); prs.slide_height=Emu(6858000)
BLANK=prs.slide_layouts[6]
RUNNING="EXPATRIATION & EXIT TAX · 2026"
FOOTER="ANTOINE COMAR · RÉUNION DE CADRAGE"

def slide(dark=False):
    s=prs.slides.add_slide(BLANK)
    if dark:
        r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,prs.slide_height)
        r.fill.solid(); r.fill.fore_color.rgb=BLACK; r.line.fill.background(); r.shadow.inherit=False
    return s

def box(s,l,t,w,h,paras,size,bold=False,color=INK,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,space=6,line_sp=None):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    if isinstance(paras,str): paras=[paras]
    for i,ptext in enumerate(paras):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(space)
        if line_sp: p.line_spacing=line_sp
        run=p.add_run(); run.text=ptext
        run.font.name=F; run.font.size=Pt(size); run.font.bold=bold; run.font.color.rgb=color
    return tb

def rule(s,l,t,w,color=RULEL,h=0.018):
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb=color; r.line.fill.background(); r.shadow.inherit=False
    return r

def chrome(s,n,total,dark=False):
    box(s,6.5,0.52,6.28,0.3,RUNNING,9,True,GOLD,PP_ALIGN.RIGHT)
    fcol=GREY9 if dark else GREYK
    box(s,0.55,7.04,8.0,0.3,FOOTER,8,False,fcol)
    box(s,11.0,7.04,1.78,0.3,f"{n:02d} / {total:02d}",8,False,fcol,PP_ALIGN.RIGHT)

def kicker(s,t,title,dark=False,title_size=28):
    sq=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(0.56),Inches(1.55),Inches(0.12),Inches(0.12))
    sq.fill.solid(); sq.fill.fore_color.rgb=GOLD; sq.line.fill.background(); sq.shadow.inherit=False
    box(s,0.8,1.42,11.8,0.32,t,11,True,WHITE if dark else INK)
    box(s,0.55,1.92,12.2,0.9,title,title_size,True,WHITE if dark else INK)
    rule(s,0.56,3.05,12.22,RULED if dark else RULEL)

def cover(n,total,title,subtitle,baseline):
    s=slide(dark=True); chrome(s,n,total,True)
    rule(s,0.56,3.5,2.6,GOLD,0.03)
    box(s,0.5,3.62,12.2,1.4,title,54,True,WHITE)
    box(s,0.56,5.25,11.0,0.6,subtitle,17,False,RGBColor(0xEC,0xEC,0xEC))
    box(s,0.56,6.0,11.5,0.4,baseline,12.5,True,WHITE)

def divider(n,total,num,title,subtitle):
    s=slide(dark=True); chrome(s,n,total,True)
    box(s,0.55,2.4,3.0,1.4,num,80,True,GOLD)
    box(s,0.58,3.95,11.0,0.7,title,34,True,WHITE)
    rule(s,0.6,4.78,2.4,GOLD,0.03)
    box(s,0.58,4.95,10.0,0.4,subtitle,15,False,GREYD)

def kpi(n,total,kick,title,items):
    s=slide(dark=True); chrome(s,n,total,True); kicker(s,kick,title,True,30)
    xs=[0.56,3.62,6.68,9.74]
    for i,(big,lab) in enumerate(items):
        x=xs[i]
        if i>0:
            r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x-0.05),Inches(3.45),Inches(0.012),Inches(2.3))
            r.fill.solid(); r.fill.fore_color.rgb=RULED; r.line.fill.background(); r.shadow.inherit=False
        box(s,x,3.6,2.95,1.15,big,34,True,WHITE)
        box(s,x,4.82,2.85,0.9,lab,11.5,False,GREY9)
    rule(s,0.56,5.95,12.22,RULED)

def twocol(n,total,kick,title,lh,lc,lb,rh,rc,rb):
    s=slide(); chrome(s,n,total); kicker(s,kick,title)
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(6.66),Inches(3.25),Inches(0.012),Inches(3.0))
    r.fill.solid(); r.fill.fore_color.rgb=RULEL; r.line.fill.background(); r.shadow.inherit=False
    box(s,0.56,3.3,5.8,0.3,lh,11,True,lc)
    box(s,0.56,3.75,5.7,2.6,lb,12.5,False,INK,space=8,line_sp=1.04)
    box(s,6.95,3.3,5.8,0.3,rh,11,True,rc)
    box(s,6.95,3.75,5.7,2.6,rb,12.5,False,INK,space=8,line_sp=1.04)

def content(n,total,kick,title,intro,body,title_size=27):
    s=slide(); chrome(s,n,total); kicker(s,kick,title,title_size=title_size)
    y=3.3
    if intro:
        box(s,0.56,3.25,12.0,0.6,intro,12.5,False,GREYK); y=3.95
    box(s,0.56,y,12.0,6.7-y,body,12.5,False,INK,space=8,line_sp=1.04)

def threecol(n,total,kick,title,intro,cols):
    s=slide(dark=True); chrome(s,n,total,True); kicker(s,kick,title,True,30)
    box(s,0.56,3.05,11.6,0.6,intro,12.5,False,GREYD)
    rule(s,0.56,3.95,12.22,RULED); xs=[0.56,4.64,8.72]
    for i,(num,hdr,bd) in enumerate(cols):
        x=xs[i]
        if i>0:
            r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x-0.12),Inches(4.2),Inches(0.012),Inches(2.2))
            r.fill.solid(); r.fill.fore_color.rgb=RULED; r.line.fill.background(); r.shadow.inherit=False
        box(s,x,4.2,3.7,0.4,num,13,True,GOLD)
        box(s,x,4.6,3.7,0.4,hdr,15,True,WHITE)
        box(s,x,5.08,3.7,1.5,bd,11.5,False,GREY9,space=4,line_sp=1.03)

def numbered(n,total,kick,title,items,note=None):
    s=slide(); chrome(s,n,total); kicker(s,kick,title)
    rows=[(0.56,3.4),(6.7,3.4),(0.56,5.0),(6.7,5.0)]
    for i,(num,hd,bd) in enumerate(items):
        x,y=rows[i]
        box(s,x,y,0.8,0.5,num,20,True,GOLD)
        box(s,x+0.75,y+0.02,5.2,0.4,hd,14,True,INK)
        box(s,x+0.75,y+0.45,5.2,1.0,bd,11.5,False,GREYK,space=3,line_sp=1.03)
    if note: box(s,0.56,6.55,12.0,0.4,note,12,True,INK)

def table(n,total,kick,title,headers,rows,note=None):
    s=slide(); chrome(s,n,total); kicker(s,kick,title)
    nr=len(rows)+1; nc=len(headers)
    gt=s.shapes.add_table(nr,nc,Inches(0.56),Inches(3.25),Inches(12.22),Inches(0.4*nr)).table
    for j,h in enumerate(headers):
        c=gt.cell(0,j); c.text=h; c.fill.solid(); c.fill.fore_color.rgb=INK
        for p in c.text_frame.paragraphs:
            for r in p.runs: r.font.name=F; r.font.size=Pt(10.5); r.font.bold=True; r.font.color.rgb=WHITE
    for i,row in enumerate(rows,1):
        for j,val in enumerate(row):
            c=gt.cell(i,j); c.text=val
            c.fill.solid(); c.fill.fore_color.rgb=WHITE if i%2 else RGBColor(0xF4,0xF4,0xF4)
            for p in c.text_frame.paragraphs:
                for r in p.runs:
                    r.font.name=F; r.font.size=Pt(10); r.font.color.rgb=INK
                    if j==0: r.font.bold=True
    if note: box(s,0.56,6.5,12.22,0.4,note,11,True,INK)

def closing(n,total,kick,title,date,signature,locations,disclaimer):
    s=slide(dark=True); chrome(s,n,total,True)
    p=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(9.55),0,Inches(3.23),prs.slide_height)
    p.fill.solid(); p.fill.fore_color.rgb=PANEL; p.line.fill.background(); p.shadow.inherit=False
    e=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(9.55),0,Inches(0.04),prs.slide_height)
    e.fill.solid(); e.fill.fore_color.rgb=GOLD; e.line.fill.background(); e.shadow.inherit=False
    box(s,0.56,1.6,8.0,0.3,kick,11,True,GOLD)
    box(s,0.52,2.05,8.7,1.7,title,38,True,WHITE)
    box(s,0.56,3.95,8.6,0.4,date,14,False,GREYD)
    rule(s,0.56,4.7,8.4,RULED,0.02)
    box(s,0.56,4.85,8.4,0.8,signature,13,True,WHITE)
    box(s,0.56,6.0,8.6,0.3,locations,9,False,GREYK)
    box(s,0.56,6.5,8.6,0.55,disclaimer,8.5,False,RGBColor(0x4A,0x4A,0x4A))

# ===================== CONTENU COMAR =====================
T=18
cover(1,T,"Expatriation & exit tax",
      "Une fenetre fiscale en 2026 pour sortir vos plus-values latentes du filet francais.",
      "ANTOINE COMAR   ·   Réunion de cadrage stratégique, 30 juin 2026")
divider(2,T,"01","LE DOSSIER","Votre projet, l'enjeu, la fenetre.")
kpi(3,T,"VOTRE SITUATION EN CHIFFRES","Ce qui est en jeu.",
    [("24/08/2026","date de départ prévue"),("≈ 3,8 M€","plus-values latentes (titres cotés)"),
     ("6 ans","le seuil de résidence à ne pas atteindre"),("> 10 M€","patrimoine familial")])
twocol(4,T,"LE POINT DÉCISIF","Tout se joue sur la durée de résidence.",
  "EN PARTANT EN 2026",GREEN,
  ["L'exit tax ne vise que ceux domiciliés en France au moins 6 des 10 années précédant le départ (art. 167 bis CGI).",
   "Rentré de Malaisie en 2020-2021, vous restez sous ce seuil au 24/08/2026.",
   "Conséquence : vous devriez être hors du champ. Les ≈ 3,8 M€ de plus-values latentes ne sont pas imposés au départ."],
  "LE SEUIL À NE PAS FRANCHIR",RED,
  ["Au-delà de 6 ans de résidence (selon la date de retour retenue), vous entrez dans le champ.",
   "L'exit tax frappe alors les ≈ 3,8 M€ de plus-values latentes.",
   "La marge est étroite : il faut sécuriser la date de retour et arrêter la date de départ."])
divider(5,T,"02","LA STRATÉGIE","Verrouiller la fenetre, choisir le pays.")
content(6,T,"LA RÈGLE","La condition des six ans, précisément.",
  "Le transfert du domicile hors de France peut déclencher l'imposition des plus-values latentes sur les titres (art. 167 bis CGI).",
  ["Condition : avoir été domicilié en France (art. 4 B), imposable sur l'ensemble de ses revenus, au moins 6 des 10 années précédant le transfert (BOFiP BOI-RPPM-PVBMI-50-10-10-10).",
   "Le décompte se fait de date à date, à partir de la date de transfert ; la condition peut être remplie de façon continue ou discontinue.",
   "Les années incomplètes comptent au prorata des jours : on additionne les jours de résidence, le total devant atteindre 6 années civiles.",
   "Appréciation par personne. Seuils du champ : portefeuille > 800 000 € ou participation >= 50 %."])
content(7,T,"LE POINT À VERROUILLER","Votre date de retour : le foyer prime.",
  "Tout dépend de la date à laquelle vous êtes redevenu résident français. Le critère du foyer (art. 4 B) commande.",
  ["Le foyer est le lieu où vous habitez normalement et avez le centre de votre vie familiale ; il prime sur le séjour principal (CE 11 mai 2022 n°450692 ; CE 27 janv. 2010 n°294784).",
   "Votre épouse et vos enfants étant rentrés dès septembre 2020, l'administration pourrait fixer votre foyer, donc votre résidence, à cette date, même en travaillant encore en Malaisie (CE 9 juin 2021 n°431551 ; CE 26 sept. 2012 n°346556).",
   "À l'inverse, la convention franco-malaisienne peut vous attribuer la résidence en Malaisie jusqu'au 06/03/2021 (foyer et activité sur place), repoussant le point de départ.",
   "Cet arbitrage de date est le coeur du dossier : il se documente."])
content(8,T,"LE CALCUL, AU JOUR PRÈS","Confortable, ou tout juste sous le seuil.",
  "Selon la date de retour retenue, vous passez largement, ou de peu, sous les 6 ans au 24/08/2026.",
  ["Si la résidence est rétablie en septembre 2020 (retour du foyer) : environ 5 ans et 11 mois, soit un peu moins de 6 ans. La marge se compte en semaines.",
   "Si elle n'est rétablie qu'en mars 2021 (résident de Malaisie au sens de la convention jusque-là) : environ 5 ans et demi, marge confortable.",
   "Deux leviers de sécurisation : documenter la résidence malaisienne jusqu'en mars 2021 ; au besoin, avancer légèrement la date de départ.",
   "Objectif : bâtir une marge nette, pour ne pas dépendre d'un décompte au jour près en cas de contrôle."])
content(9,T,"LE FILET DE SÉCURITÉ","Si l'exit tax s'appliquait malgré tout.",
  "Même si la condition de durée était regardée comme remplie, le départ reste gérable.",
  ["Sursis de plein droit, sans garantie, vers l'UE / l'EEE ou un État lié à la France par convention d'assistance et de recouvrement, hors État non coopératif (art. 167 bis IV).",
   "Vers les autres États, dont le Panama (non coopératif, art. 238-0 A) : sursis possible mais SUR DEMANDE, avec représentant fiscal et garanties de 12,8 % du montant brut des plus-values, à constituer avant le départ (art. 167 bis V).",
   "Dégrèvement à l'expiration du délai légal si les titres sont conservés.",
   "La convention d'établissement franco-panaméenne ne dispense pas de ces conditions (Cass. 1994 et 1999)."])
table(10,T,"COMPARATIF","Trois destinations envisagées.",
  ["","Île Maurice","Paraguay","Panama"],
  [["Convention fiscale FR","Oui","Non","Non"],
   ["Sursis exit tax (si applicable)","De plein droit (à confirmer)","Sur demande + garanties","Sur demande + garanties (ETNC)"],
   ["Dividendes / plus-values","Régime attractif","Territorial","Territorial"],
   ["Point de vigilance","Substance réelle","Pas de convention","Listes / image"],
   ["Lecture cabinet","Favorable","À sécuriser","À sécuriser"]],
  "Règle constante du cabinet : pas de destination sans convention sans sécurisation préalable.")
content(11,T,"APRÈS LE DÉPART","Votre fiscalité une fois non-résident.",
  "Sortir du champ de l'exit tax ne règle pas tout : la résidence d'arrivée gouverne vos revenus futurs.",
  ["Plus-values sur titres étrangers (portefeuille IBKR) : hors champ français une fois non-résident, donc taxées selon le pays d'accueil.",
   "Exception (art. 244 bis B) : la cession de titres d'une société française détenue à plus de 25 % reste imposable en France ; vos holdings françaises (Sifa, Sofidil, SCI) sont à cartographier sous cet angle.",
   "Dividendes (≈ 288 k€/an) : imposés selon la résidence et les conventions ; retenues à la source à anticiper.",
   "Revenus de source française (SCI, fonciers ≈ 27 k€) : restent imposables en France. Le choix du pays se joue sur les revenus passifs, pas seulement sur l'exit tax."])
content(12,T,"CALENDRIER","La séquence des opérations.",
  "L'ordre et la date des cessions sont déterminants ; ils se décident maintenant.",
  ["Déjà réalisé, en tant que résident (taxé en France) : plus-values 2025 de 694 k€ ; plus-values 2026 d'environ 940 k€.",
   "Avant le départ : céder les titres en moins-value (≈ 225 k€) pour imputer sur les plus-values 2026.",
   "Départ le 24/08/2026, dans la fenetre des moins de 6 ans.",
   "Après le départ : céder le portefeuille en plus-value une fois non-résident, hors du champ de l'exit tax."])
content(13,T,"À SÉCURISER EN AMONT","Les sujets à traiter avant de partir.",None,
  ["Régularisation des comptes étrangers (Interactive Brokers, HSBC Singapour) : à finaliser avant le départ (≈ 44 k€ hors pénalités ; amende art. 1736 IV, atténuée en cas de régularisation spontanée).",
   "Holding AC HOLDING 1 (radiée 02/2025) : vérifier l'absence de plus-value en report (apport-cession, art. 150-0 B ter) que le départ remobiliserait dans l'exit tax.",
   "Votre épouse Sibylle : résidence et situation propre (nue-propriété, micro-entreprise) à traiter en parallèle.",
   "Substance dans le pays d'accueil (logement, présence, centre des intérêts) et scolarité des trois enfants, à aligner sur le calendrier fiscal."])
divider(14,T,"03","NOTRE MISSION","Confirmer, sécuriser, mettre en oeuvre.")
threecol(15,T,"NOTRE ACCOMPAGNEMENT","Nous sécurisons votre départ de bout en bout.",
  "De la confirmation de la fenetre des 6 ans jusqu'à l'installation, nous confirmons, sécurisons et pilotons.",
  [("01","CONFIRMER","Établir, jour par jour et pièces à l'appui, que la domiciliation française reste sous 6 ans, et fixer la date de départ."),
   ("02","SÉCURISER","Régularisation des comptes, choix du pays, calendrier des cessions et de la résidence."),
   ("03","METTRE EN OEUVRE","Formalités de départ, preuve de résidence et suivi post-départ, pour vous et votre épouse.")])
content(16,T,"NOTRE PROPOSITION","Un cadrage, puis un forfait global.",None,
  ["ÉTAPE 1, le cadrage stratégique : analyse de votre situation, confirmation de la fenetre, comparatif des destinations, calendrier et points de vigilance. Livrable : un document de travail opérationnel. Honoraires : 2 400 € HT.",
   "ÉTAPE 2, l'accompagnement au forfait : défini après le cadrage, il couvre la mise en oeuvre jusqu'au départ et le suivi post-départ.",
   "Au regard de l'enjeu (≈ 3,8 M€ de plus-values en jeu), l'accompagnement se raisonne comme une assurance du dispositif."])
numbered(17,T,"ET MAINTENANT","Nos premières actions.",
  [("01","Figer la date de retour","Réunir les preuves de résidence 2020-2021, et de la résidence malaisienne jusqu'en mars 2021."),
   ("02","Arrêter la date de départ","Avant le seuil des 6 ans, avec une marge de sécurité."),
   ("03","Finaliser la régularisation","Comptes étrangers soldés et déclarés avant le départ."),
   ("04","Arbitrer la destination","Île Maurice, Paraguay ou Panama, au vu de vos priorités.")],
  "Pièces utiles d'ici le RDV : prix de revient du portefeuille · justificatifs de retour 2020-2021 · statuts des sociétés et SCI.")
closing(18,T,"PROCHAINE ÉTAPE","Sécurisons votre départ.",
  "Mardi 30 juin 2026, 15h00 (Paris) / 9h00 (Martinique), en visioconférence.",
  "Me François Ouairy\nAVOCAT, BENSAID AVOCATS",
  "PARIS · GENÈVE · MARSEILLE · CANNES · LISBONNE",
  "Document de travail confidentiel. Sources : CGI art. 167 bis, 4 A, 4 B, 244 bis B, 150-0 B ter, 238-0 A, 1736 IV ; BOFiP BOI-RPPM-PVBMI-50-10-10-10, BOI-INT-DG-20-10-10 ; CE 11 mai 2022 n°450692, 9 juin 2021 n°431551, 21 nov. 2012 n°347223 ; CJCE 11 mars 2004 C-9/02. À confirmer au vu des pièces complètes.")

out=sys.argv[1] if len(sys.argv)>1 else "/tmp/support_comar.pptx"
prs.save(out)
print("OK slides:",len(prs.slides._sldIdLst),"->",out)
